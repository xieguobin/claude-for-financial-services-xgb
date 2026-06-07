#!/usr/bin/env python3
"""
Generic A-share Pitch Deck Generator v2
========================================
Data sources (all work without proxy):
  - Tencent qt.gtimg.cn   → real-time quote, peer batch, K-line history
  - 同花顺 THS (AkShare)  → annual financial abstracts
Usage:
    python3 generate_a_share_ppt.py --company "宁德时代" --ticker "300750" --industry "电池"
    python3 generate_a_share_ppt.py --company "贵州茅台" --ticker "600519" --industry "白酒"
"""
import os, sys, json, argparse, re, io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd

# ─── Proxy bypass: set before any network imports ───
os.environ['NO_PROXY'] = '*'
os.environ['no_proxy'] = '*'
for _k in ('https_proxy', 'http_proxy', 'all_proxy',
           'HTTPS_PROXY', 'HTTP_PROXY', 'ALL_PROXY'):
    os.environ.pop(_k, None)

import requests as _req_lib

# Patch Session so every new session has trust_env=False
_SESSION_CLASS = _req_lib.Session
_orig_session_init = _SESSION_CLASS.__init__
def _patched_session_init(self, *a, **kw):
    _orig_session_init(self, *a, **kw)
    self.trust_env = False
_SESSION_CLASS.__init__ = _patched_session_init

# ─── Color palette ───
C_BG     = RGBColor(0xFF, 0xFF, 0xFF)
C_PRI    = RGBColor(0x1A, 0x3C, 0x6E)
C_ACCENT = RGBColor(0xC8, 0xA0, 0x32)
C_TEXT   = RGBColor(0x33, 0x33, 0x33)
C_LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x66, 0x66, 0x66)
C_DIM    = RGBColor(0x99, 0x99, 0x99)

CHART_DIR = "/tmp"
SESSION = _SESSION_CLASS()

# ═══════════════════════════════════════════════════════════
#  DATA LAYER  — 全部从公开API实时获取
# ═══════════════════════════════════════════════════════════

def _get(url, timeout=15):
    return SESSION.get(url, timeout=timeout)

def _prefix(ticker):
    """SH for 6xx/9xx, SZ for 0xx/3xx/2xx"""
    return "sh" if ticker[0] in ('6', '9') else "sz"

# ── 1. 实时行情 (Tencent) ──
def fetch_quote(ticker):
    prefix = _prefix(ticker)
    r = _get(f"https://qt.gtimg.cn/q={prefix}{ticker}")
    text = r.text.strip()
    # Delimiter varies by market: SZ stocks use =" while SH stocks use ="~
    sep = '="~' if '="~' in text else '="'
    if sep not in text:
        return None
    parts = text.split(sep)[1].rstrip('";').split('~')
    if len(parts) < 5:
        return None
    # Fields vary by market; safely extract
    def _f(idx):
        return parts[idx] if idx < len(parts) and parts[idx] else None
    result = {'name': _f(1).replace(' ', ''), 'code': _f(2),
              'price': _f(3), 'prev_close': _f(4), 'open': _f(5),
              'high': _f(33), 'low': _f(34),
              'volume': _f(36), 'amount': _f(37), 'turnover': _f(38),
              'pe': _f(39), 'pb': _f(46), 'mcap': _f(45)}
    # Convert numerics
    for k in ('price', 'prev_close', 'open', 'high', 'low', 'pe', 'pb', 'mcap'):
        v = result[k]
        result[k] = float(v) if v else None
    if result['mcap']:
        result['mcap_yi'] = result['mcap'] / 10000  # 元→亿元
    return result

# ── 2. 财务摘要 (同花顺 via AkShare) ──
def fetch_financials(ticker):
    import akshare as ak
    try:
        df = ak.stock_financial_abstract_ths(symbol=ticker, indicator="年度")
        if df is not None and not df.empty:
            # Dates may be "2024" (year only) or "2024-12-31"
            yr_col = df.iloc[:, 0].astype(str)
            # Filter to annual rows: either year-only >= 2000 or contain date pattern
            mask = yr_col.str.match(r'^\d{4}$') | yr_col.str.contains(r'\d{4}-\d{2}-\d{2}')
            annual = df[mask]
            # Take last 5 years
            return annual.tail(5).reset_index(drop=True)
    except Exception as e:
        print(f"  [warn] 财务数据获取失败: {e}")
    return pd.DataFrame()

# ── 3. 历史K线 (Tencent) ──
def fetch_kline(ticker, start="20240101", end="20250608"):
    prefix = _prefix(ticker)
    # Tencent K-line API requires YYYY-MM-DD format
    def _fmt(d):
        d = str(d)
        if len(d) == 8 and '-' not in d:
            return f"{d[:4]}-{d[4:6]}-{d[6:8]}"
        return d
    start_f, end_f = _fmt(start), _fmt(end)
    url = (f"https://web.ifzq.gtimg.cn/appstock/app/fqkline/get"
           f"?param={prefix}{ticker},day,{start_f},{end_f},500,qfq")
    r = _get(url, timeout=15)
    data = r.json()
    stock_key = f"{prefix}{ticker}"
    stock_data = data.get('data', {}).get(stock_key, {})
    if isinstance(stock_data, dict):
        klines = stock_data.get('qfqday', []) or stock_data.get('day', [])
    elif isinstance(data.get('data'), list):
        klines = data['data']
    else:
        klines = []
    if not klines:
        return pd.DataFrame()
    # Strip dividend info dict (7th element) if present
    cleaned = [row[:6] for row in klines]
    cols = ['日期', '开盘', '收盘', '最高', '最低', '成交量']
    df = pd.DataFrame(cleaned, columns=cols)
    for col in cols[1:]:
        df[col] = pd.to_numeric(df[col], errors='coerce')
    df['日期'] = pd.to_datetime(df['日期'])
    return df
    for col in cols[1:]:
        df[col] = pd.to_numeric(df[col], errors='coerce')
    df['日期'] = pd.to_datetime(df['日期'])
    return df

# ── 4. 同行可比公司 (Tencent 批量) ──
def fetch_peers(tickers):
    """tickers: list of 6-digit codes. Returns list of quote dicts."""
    if not tickers:
        return []
    pairs = [f"{_prefix(t)}{t}" for t in tickers]
    r = _get(f"https://qt.gtimg.cn/q={','.join(pairs)}")
    results = []
    for line in r.text.strip().split(';'):
        line = line.strip()
        if not line:
            continue
        sep = '="~' if '="~' in line else '="'
        if sep not in line:
            continue
        parts = line.split(sep)[1].rstrip('";').split('~')
        if len(parts) < 4 or not parts[1]:
            continue
        def _f(idx):
            return parts[idx] if idx < len(parts) and parts[idx] else None
        pe_str = _f(39)
        pb_str = _f(46)
        mcap_str = _f(45)
        pe = float(pe_str) if pe_str and 0 < float(pe_str) < 999 else None
        pb = float(pb_str) if pb_str and float(pb_str) > 0 else None
        mcap = float(mcap_str) / 10000 if mcap_str else None
        results.append({
            'name': _f(1).replace(' ', ''),
            'code': _f(2),
            'price': float(_f(3)) if _f(3) else None,
            'pe': pe, 'pb': pb, 'mcap_yi': mcap,
        })
    return results


# ═══════════════════════════════════════════════════════════
#  CHART GENERATION
# ═══════════════════════════════════════════════════════════

def _save(fig, name):
    path = os.path.join(CHART_DIR, name)
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def _to_float(series):
    """Strip 亿/% etc and convert to float."""
    return series.astype(str).str.replace('亿', '').str.replace('%', '').str.replace(',', '').replace('False', '').replace('nan', '').astype(float)

def _f(val):
    """Convert a single value: strip 亿/% and return float or None."""
    if val is None:
        return None
    s = str(val).replace('亿', '').replace('%', '').replace(',', '').strip()
    if not s or s in ('False', 'nan', 'None'):
        return None
    try:
        return float(s)
    except:
        return None

def make_charts(fin, hist, peers, company, ticker, chart_dir):
    paths = {}
    years = fin['报告期'].astype(str).tolist() if not fin.empty and '报告期' in fin.columns else []
    rev = _to_float(fin['营业总收入']).tolist() if not fin.empty and '营业总收入' in fin.columns else []
    profit = _to_float(fin['净利润']).tolist() if not fin.empty and '净利润' in fin.columns else []

    # 1) Revenue & Profit
    if rev and profit:
        fig, ax = plt.subplots(figsize=(7, 3.5))
        x, w = np.arange(len(years)), 0.35
        b1 = ax.bar(x - w/2, rev, w, label='营业收入', color='#1A3C6E', alpha=0.85)
        b2 = ax.bar(x + w/2, profit, w, label='归母净利润', color='#C8A032', alpha=0.85)
        ax.set_xticks(x); ax.set_xticklabels(years)
        ax.set_ylabel('亿元'); ax.set_title(f'{company} — 营业收入与归母净利润', fontsize=12, fontweight='bold', color='#1A3C6E', pad=10)
        ax.legend(loc='upper left'); ax.grid(True, alpha=0.3, axis='y')
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        for b in b1: ax.text(b.get_x()+b.get_width()/2, b.get_height()*1.02, f'{b.get_height():.0f}', ha='center', fontsize=8)
        for b in b2: ax.text(b.get_x()+b.get_width()/2, b.get_height()*1.02, f'{b.get_height():.0f}', ha='center', fontsize=8)
        fig.tight_layout(); paths['rev_profit'] = _save(fig, f'{ticker}_01_rev_profit.png')

    # 2) Margins
    gm = _to_float(fin['销售毛利率']).tolist() if not fin.empty and '销售毛利率' in fin.columns else []
    nm = _to_float(fin['销售净利率']).tolist() if not fin.empty and '销售净利率' in fin.columns else []
    roe = _to_float(fin['净资产收益率']).tolist() if not fin.empty and '净资产收益率' in fin.columns else []
    if gm:
        fig, ax = plt.subplots(figsize=(7, 3.5))
        if gm: ax.plot(years, gm, 'o-', color='#1A3C6E', lw=2, ms=6, label='毛利率')
        if nm: ax.plot(years, nm, 's-', color='#C8A032', lw=2, ms=6, label='净利率')
        if roe: ax.plot(years, roe, '^-', color='#E67E22', lw=2, ms=6, label='ROE')
        ax.set_ylabel('%'); ax.set_title(f'{company} — 盈利能力指标', fontsize=12, fontweight='bold', color='#1A3C6E', pad=10)
        ax.legend(loc='lower right'); ax.grid(True, alpha=0.3)
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        fig.tight_layout(); paths['margins'] = _save(fig, f'{ticker}_02_margins.png')

    # 3) Growth rate
    if len(rev) >= 2:
        rg = [None] + [round((rev[i]-rev[i-1])/rev[i-1]*100, 1) for i in range(1, len(rev))]
        pg = [None] + [round((profit[i]-profit[i-1])/profit[i-1]*100, 1) for i in range(1, len(profit))]
        fig, ax = plt.subplots(figsize=(7, 3))
        x = np.arange(len(years))
        ax.bar(x-0.2, [g or 0 for g in rg], 0.4, label='营收增长率', color='#1A3C6E', alpha=0.85)
        ax.bar(x+0.2, [g or 0 for g in pg], 0.4, label='净利润增长率', color='#C8A032', alpha=0.85)
        ax.set_xticks(x); ax.set_xticklabels(years); ax.set_ylabel('%')
        ax.set_title(f'{company} — 同比增长率', fontsize=12, fontweight='bold', color='#1A3C6E', pad=10)
        ax.legend(); ax.grid(True, alpha=0.3, axis='y')
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        for i, (a, b) in enumerate(zip(rg, pg)):
            if a: ax.text(i-0.2, a+0.5, f'{a:.1f}%', ha='center', fontsize=8)
            if b: ax.text(i+0.2, b+0.5, f'{b:.1f}%', ha='center', fontsize=8)
        fig.tight_layout(); paths['growth'] = _save(fig, f'{ticker}_03_growth.png')

    # 4) Peer comparison
    if peers:
        # Filter to peers with valid price AND pe/pb for charts
        valid_peers = [p for p in peers if p.get('price') and p.get('pe') is not None and p.get('pb') is not None]
        pnames = [p['name'] for p in valid_peers]
        ppe = [p['pe'] for p in valid_peers]
        ppb = [p['pb'] for p in valid_peers]
        if pnames:
            fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(8, 3.5))
            c1 = ['#C8A032' if ticker in (p.get('code','')) else '#1A3C6E' for p in valid_peers]
            ax1.barh(pnames, ppe, color=c1, alpha=0.85)
            ax1.set_xlabel('P/E'); ax1.set_title('可比公司 P/E', fontsize=11, fontweight='bold', color='#1A3C6E')
            ax1.grid(True, alpha=0.3, axis='x'); ax1.spines['top'].set_visible(False); ax1.spines['right'].set_visible(False)
            for bar, val in zip(ax1.patches, ppe):
                if val: ax1.text(val+0.5, bar.get_y()+bar.get_height()/2, f'{val:.1f}x', va='center', fontsize=9)
            c2 = ['#C8A032' if ticker in (p.get('code','')) else '#1A3C6E' for p in valid_peers]
            ax2.barh(pnames, ppb, color=c2, alpha=0.85)
            ax2.set_xlabel('P/B'); ax2.set_title('可比公司 P/B', fontsize=11, fontweight='bold', color='#1A3C6E')
            ax2.grid(True, alpha=0.3, axis='x'); ax2.spines['top'].set_visible(False); ax2.spines['right'].set_visible(False)
            for bar, val in zip(ax2.patches, ppb):
                if val: ax2.text(val+0.1, bar.get_y()+bar.get_height()/2, f'{val:.1f}x', va='center', fontsize=9)
            fig.tight_layout(); paths['peers'] = _save(fig, f'{ticker}_04_peers.png')

    # 5) Price trend
    if hist is not None and not hist.empty:
        fig, ax = plt.subplots(figsize=(8, 3.5))
        ax.fill_between(hist['日期'], hist['收盘'], alpha=0.15, color='#1A3C6E')
        ax.plot(hist['日期'], hist['收盘'], color='#1A3C6E', lw=2)
        ax.axhline(y=hist['收盘'].mean(), color='#C8A032', ls='--', alpha=0.5, label=f'均价: ¥{hist["收盘"].mean():.0f}')
        ax.set_title(f'{company} ({ticker}) — 股价走势', fontsize=13, fontweight='bold', color='#1A3C6E', pad=10)
        ax.set_ylabel('¥'); ax.legend(loc='upper right'); ax.grid(True, alpha=0.3)
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        fig.tight_layout(); paths['price'] = _save(fig, f'{ticker}_05_price.png')

    # 6) Valuation football field
    if peers and quote and quote.get('pe'):
        valid_pes = [p['pe'] for p in peers if p['pe'] and 0 < p['pe'] < 500]
        if valid_pes:
            avg_pe = np.mean(valid_pes)
            cp, cpe = quote['price'], quote['pe']
            pe_low = cp * avg_pe * 0.85 / cpe
            pe_mid = cp * avg_pe / cpe
            pe_high = cp * avg_pe * 1.18 / cpe
            methods = ['行业平均PE', 'PE下限(-15%)', 'PE中枢', 'PE上限(+18%)']
            vals = [cp * avg_pe / cpe, pe_low, pe_mid, pe_high]
            colors = ['#2980B9', '#1A3C6E', '#C8A032', '#1A3C6E']
            fig, ax = plt.subplots(figsize=(8, 3))
            ax.barh(methods, vals, color=colors, alpha=0.8, height=0.6)
            ax.axvline(x=cp, color='#C0392B', lw=2.5, ls='--', label=f'当前价: ¥{cp:.2f}')
            ax.set_xlabel('¥'); ax.set_title(f'{company} — 估值区间', fontsize=12, fontweight='bold', color='#1A3C6E', pad=10)
            ax.legend(loc='lower right'); ax.grid(True, alpha=0.3, axis='x')
            ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
            for bar, val in zip(ax.patches, vals):
                ax.text(val+20, bar.get_y()+bar.get_height()/2, f'¥{val:.0f}', va='center', fontsize=9)
            fig.tight_layout(); paths['valuation'] = _save(fig, f'{ticker}_06_val.png')

    return paths


# ═══════════════════════════════════════════════════════════
#  PPT BUILDER
# ═══════════════════════════════════════════════════════════

def _add_header(slide, text):
    s = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(0.8))
    s.fill.solid(); s.fill.fore_color.rgb = C_PRI; s.line.color.rgb = C_PRI
    tf = s.text_frame
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = C_WHITE; p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5); tf.margin_top = Inches(0.15)

def _add_text(slide, text, l, t, w, h, size=13, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.alignment = align

def _add_bullets(slide, items, l, t, w, h, size=13, color=C_TEXT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = color; p.space_after = Pt(6)

def _stat_box(slide, items, x, y, cols=3, cw=4.2, rh=0.55):
    for i, (label, value, note) in enumerate(items):
        col, row = i % cols, i // cols
        s = slide.shapes.add_shape(1, Inches(x+col*cw), Inches(y+row*rh), Inches(cw-0.15), Inches(rh-0.05))
        s.fill.solid(); s.fill.fore_color.rgb = C_LIGHT if i%2==0 else C_BG
        s.line.color.rgb = RGBColor(0xDD,0xDD,0xDD)
        tf = s.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"{label}  {value}  {('('+note+')') if note else ''}"
        p.font.size = Pt(10.5); p.font.color.rgb = C_TEXT


def build_deck(args, quote, fin, hist, peers, charts):
    co, tk = args.company, args.ticker
    ex = "SH" if tk[0] in ('6','9') else "SZ"
    years = fin['报告期'].tolist() if not fin.empty and '报告期' in fin.columns else []
    today = pd.Timestamp.now().strftime('%Y年%m月%d日')
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def blank():
        return prs.slides.add_slide(prs.slide_layouts[6])

    # ── S1: Cover ──
    s = blank()
    bg = s.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = C_PRI; bg.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = co; p.font.size = Pt(52); p.font.bold = True; p.font.color.rgb = C_WHITE; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = f"{tk}.{ex}  |  {today}"; p2.font.size = Pt(17); p2.font.color.rgb = C_ACCENT; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(12)
    p3 = tf.add_paragraph(); p3.text = args.title or "投资价值分析  |  深度路演材料"; p3.font.size = Pt(22); p3.font.color.rgb = C_WHITE; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(18)
    _add_text(s, f"机密文件 | 仅供内部参考 | 数据来源: 腾讯财经 / 同花顺THS / AkShare | {today}",
              Inches(1), Inches(6.6), Inches(11), Inches(0.4), size=10, color=RGBColor(0xCC,0xCC,0xCC), align=PP_ALIGN.CENTER)

    # ── S2: Highlights ──
    s = blank(); _add_header(s, f"投资摘要  |  Investment Highlights — {co}")
    def _sf(q, k, fmt=''):
        v = q.get(k) if q else None
        return fmt.format(v) if v else "N/A"
    last = fin.iloc[-1] if not fin.empty else None
    metrics = [
        ("当前股价", _sf(quote,'price','¥{:.2f}'), ''),
        ("总市值", _sf(quote,'mcap_yi','{:.0f}亿'), ''),
        ("动态市盈率", _sf(quote,'pe','{:.1f}x'), ''),
        ("市净率", _sf(quote,'pb','{:.2f}x'), ''),
        ("最新营收", f"{_f(last['营业总收入']):.0f}亿" if last is not None and '营业总收入' in last else "N/A", ''),
        ("最新净利润", f"{_f(last['净利润']):.0f}亿" if last is not None and '净利润' in last else "N/A", ''),
    ]
    for i, (lb, val, _) in enumerate(metrics):
        col, row = i % 3, i // 3
        s2 = s.shapes.add_shape(1, Inches(0.5+col*4.2), Inches(1.1+row*0.8), Inches(3.9), Inches(0.65))
        s2.fill.solid(); s2.fill.fore_color.rgb = C_LIGHT; s2.line.color.rgb = C_ACCENT
        tf = s2.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = val; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = C_PRI; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = lb; p2.font.size = Pt(9); p2.font.color.rgb = C_GRAY; p2.alignment = PP_ALIGN.CENTER

    highlights = []
    if not fin.empty:
        highlights.append(f"财务数据：已获取 {len(fin)} 年历史财务摘要（同花顺THS）")
        try:
            gm = _f(last['销售毛利率']); nm = _f(last['销售净利率']); roe = _f(last['净资产收益率'])
            highlights.append(f"盈利能力：毛利率 {gm:.1f}%，净利率 {nm:.1f}%，ROE {roe:.1f}%")
        except: pass
    if quote:
        highlights.append(f"实时报价：¥{quote['price']}  PE(TTM) {quote.get('pe','N/A')}  PB {quote.get('pb','N/A')}")
    if peers:
        highlights.append(f"行业对标：{len(peers)} 家可比公司估值对比完成")
    highlights.append(f"数据来源：实时行情 + 历史K线({len(hist)}条) + 财务摘要 + 行业可比公司")
    while len(highlights) < 4:
        highlights.append("详见后续章节详细分析")
    _add_bullets(s, highlights[:5], Inches(0.5), Inches(2.7), Inches(12), Inches(3.8), size=14)

    # ── S3: Company Overview ──
    s = blank(); _add_header(s, f"公司概览  |  Company Overview — {co}")
    info = [f"股票代码：{tk}.{ex}"]
    if quote:
        info += [f"当前股价：¥{quote['price']}",
                 f"总市值：{quote.get('mcap_yi',0):.0f}亿元" if quote.get('mcap_yi') else "",
                 f"换手率：{quote.get('turnover')}%" if quote.get('turnover') else ""]
    info += ["", "核心财务指标（最新年度）："]
    if not fin.empty and last is not None:
        for col_cn, col_en in [('营业总收入','营业收入'),('净利润','归母净利润'),('销售毛利率','毛利率'),
                                ('销售净利率','净利率'),('净资产收益率','ROE'),('资产负债率','资产负债率'),
                                ('基本每股收益','EPS')]:
            if col_cn in last.index and pd.notna(last[col_cn]):
                try: info.append(f"  {col_en}：{_f(last[col_cn]):.2f}")
                except: info.append(f"  {col_en}：{last[col_cn]}")
    _add_bullets(s, [l for l in info if l], Inches(0.5), Inches(1.1), Inches(6), Inches(6), size=13)

    if not fin.empty and last is not None:
        _add_text(s, f"关键数据概览", Inches(6.8), Inches(1.1), Inches(6), Inches(0.4), size=14, bold=True, color=C_PRI)
        stats = []
        for lb, cn, note in [("营业收入","营业总收入",""),("归母净利润","净利润",""),
                              ("毛利率","销售毛利率",""),("净利率","销售净利率",""),
                              ("ROE","净资产收益率",""),("EPS","基本每股收益","元"),
                              ("资产负债率","资产负债率",""),("经营现金流","","充裕")]:
            if cn in last.index and pd.notna(last[cn]):
                try: v = f"{_f(last[cn]):.2f}"
                except: v = str(last[cn])
                stats.append((lb, v, note))
        if stats:
            _stat_box(s, stats, 6.8, 1.6, cols=1, cw=6, rh=0.5)

    # ── S4: Industry ──
    s = blank(); _add_header(s, f"行业分析  |  Industry Overview — {args.industry or co}")
    _add_bullets(s, [
        f"所属行业：{args.industry or '（请指定 --industry 参数）'}",
        f"行业成分股：{len(peers)} 家可比公司（详见估值分析页）",
        "",
        "市场规模与增长趋势：",
        "  • 行业整体规模及增速（基于成分股市值及财务数据推算）",
        "  • 集中度趋势：头部企业市场份额变化",
        "",
        "核心驱动因素：",
        "  • 消费升级 / 政策支持 / 技术创新（根据行业属性）",
        "  • 竞争格局与进入壁垒分析",
        "",
        f"可比公司：共 {len(peers)} 家（详见估值分析页）",
    ], Inches(0.5), Inches(1.1), Inches(12), Inches(6), size=14)

    # ── S5: Financial Summary ──
    s = blank(); _add_header(s, f"财务分析  |  Financial Summary — {co}")
    if 'rev_profit' in charts:
        s.shapes.add_picture(charts['rev_profit'], Inches(0.4), Inches(1.1), width=Inches(6.4))
    if 'margins' in charts:
        s.shapes.add_picture(charts['margins'], Inches(0.4), Inches(4.5), width=Inches(6.4))
    _add_text(s, "关键财务指标", Inches(7), Inches(1.1), Inches(6), Inches(0.4), size=14, bold=True, color=C_PRI)
    if not fin.empty and last is not None:
        ratios = []
        for lb, cn, note in [("营收增长率","营业总收入同比增长率",""),("净利润增长率","净利润同比增长率",""),
                              ("毛利率","销售毛利率",""),("净利率","销售净利率",""),
                              ("ROE","净资产收益率",""),("资产负债率","资产负债率","")]:
            if cn in last.index and pd.notna(last[cn]):
                try: v = f"{_f(last[cn]):.1f}%"
                except: v = str(last[cn])
                ratios.append((lb, v, note))
        if ratios:
            _stat_box(s, ratios, 7, 1.6, cols=1, cw=6, rh=0.5)
    if 'growth' in charts:
        s.shapes.add_picture(charts['growth'], Inches(7), Inches(4.5), width=Inches(6))

    # ── S6: Price Trend ──
    s = blank(); _add_header(s, f"股价走势  |  Stock Performance — {co} ({tk})")
    if 'price' in charts:
        s.shapes.add_picture(charts['price'], Inches(0.4), Inches(1.1), width=Inches(12.5))
        info_lines = []
        if hist is not None and not hist.empty:
            info_lines.append(f"区间: ¥{hist['最低'].min():.2f} ~ ¥{hist['最高'].max():.2f}  |  最新: ¥{hist.iloc[-1]['收盘']:.2f}")
        if quote:
            info_lines.append(f"实时: ¥{quote['price']}  PE: {quote.get('pe','N/A')}  PB: {quote.get('pb','N/A')}")
        _add_bullets(s, info_lines, Inches(0.5), Inches(4.7), Inches(12), Inches(2), size=13)
    else:
        _add_text(s, "价格数据获取中...", Inches(0.5), Inches(1.1), Inches(12), Inches(6), size=20, color=C_GRAY)

    # ── S7: Peer Comps ──
    s = blank(); _add_header(s, "可比公司估值  |  Peer Comparables")
    if 'peers' in charts:
        s.shapes.add_picture(charts['peers'], Inches(0.4), Inches(1.1), width=Inches(7.5))
    if peers:
        txt = "行业可比公司估值对比（实时）：\n\n"
        for p in peers:
            if p['price']:
                txt += f"  {p['name']} ({p['code']}): ¥{p['price']:.2f}"
                txt += f"  PE: {p['pe']:.1f}x" if p['pe'] else "  PE: N/A"
                txt += f"  PB: {p['pb']:.2f}x\n" if p['pb'] else "\n"
        txt += f"\n  → {co} 估值定位：详见估值分析页"
        _add_text(s, txt, Inches(8.1), Inches(1.1), Inches(4.9), Inches(5.5), size=12)

    # ── S8: Valuation ──
    s = blank(); _add_header(s, f"估值分析  |  Valuation Analysis — {co}")
    if 'valuation' in charts:
        s.shapes.add_picture(charts['valuation'], Inches(0.4), Inches(1.1), width=Inches(7.5))
    _add_text(s, "估值方法", Inches(8.2), Inches(1.1), Inches(4.8), Inches(0.4), size=14, bold=True, color=C_PRI)
    if peers and quote and quote.get('pe'):
        valid = [p['pe'] for p in peers if p['pe'] and 0 < p['pe'] < 500]
        if valid:
            avg = np.mean(valid); cp, cpe = quote['price'], quote['pe']
            v_low = cp * avg * 0.85 / cpe; v_mid = cp * avg / cpe; v_high = cp * avg * 1.18 / cpe
            for i, (lb, val, note) in enumerate([
                ("行业平均PE", f"{avg:.1f}x", ""),
                ("估值下限", f"¥{v_low:.0f}", "悲观"),
                ("估值中枢", f"¥{v_mid:.0f}", "基准"),
                ("估值上限", f"¥{v_high:.0f}", "乐观"),
            ]):
                top = Inches(1.6 + i*0.75)
                box = s.shapes.add_shape(1, Inches(8.2), top, Inches(4.8), Inches(0.65))
                box.fill.solid(); box.fill.fore_color.rgb = C_LIGHT; box.line.color.rgb = C_ACCENT
                tf = box.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.text = f"{lb}  {val}  {note}"
                p.font.size = Pt(12); p.font.bold = (i==2); p.font.color.rgb = C_PRI
            _add_text(s, f"当前 ¥{cp} | PE {cpe:.1f}x | PB {quote.get('pb','N/A')}",
                      Inches(8.2), Inches(4.7), Inches(4.8), Inches(0.4), size=9, color=C_DIM)

    # ── S9: Thesis ──
    s = blank(); _add_header(s, f"投资逻辑  |  Investment Thesis — {co}")
    theses = [
        ("行业地位", f"{co} 在 {args.industry or '所属行业'} 中的竞争定位与市场份额"),
        ("财务质量", f"毛利率、净利率、ROE 等核心盈利能力指标（详见财务分析页，基于 {len(fin)} 年数据）"),
        ("成长驱动", f"营收和利润增长趋势分析，基于 {len(hist)} 条历史K线数据"),
        ("估值水平", f"当前PE/PB与 {len(peers)} 家行业可比公司对比分析"),
        ("风险因素", "行业风险、公司风险、市场风险综合评估（详见风险提示页）"),
    ]
    for i, (title, desc) in enumerate(theses):
        col, row = i % 3, i // 3
        box = s.shapes.add_shape(1, Inches(0.4+col*4.3), Inches(1.1+row*2.9), Inches(4), Inches(2.6))
        box.fill.solid(); box.fill.fore_color.rgb = C_LIGHT; box.line.color.rgb = C_ACCENT
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = C_PRI; p.space_after = Pt(8)
        p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(11); p2.font.color.rgb = C_TEXT

    # ── S10: Risks ──
    s = blank(); _add_header(s, "风险提示  |  Risk Factors")
    risks = [
        ("宏观经济风险", "宏观经济波动可能影响行业需求和公司业绩表现"),
        ("行业竞争风险", "行业竞争加剧，市场份额及毛利率可能受到压力"),
        ("政策监管风险", "行业政策、环保法规变化可能对公司经营产生影响"),
        ("市场波动风险", "股票价格受市场情绪、资金面、流动性等多因素影响"),
        ("业绩不及预期", "公司实际业绩可能低于市场一致预期"),
        ("估值回调风险", "若行业景气度下行或市场风格切换，可能面临估值压缩"),
    ]
    for i, (title, desc) in enumerate(risks):
        col, row = i % 2, i // 2
        box = s.shapes.add_shape(1, Inches(0.5+col*6.3), Inches(1.1+row*1.65), Inches(6.1), Inches(1.5))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFD,0xF3,0xF3); box.line.color.rgb = C_RED
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = f"⚠️ {title}"; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = C_RED; p.space_after = Pt(5)
        p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(11); p2.font.color.rgb = C_TEXT

    # ── S11: Recommendation ──
    s = blank(); _add_header(s, f"投资建议  |  Investment Recommendation — {co}")
    box = s.shapes.add_shape(1, Inches(2), Inches(1.3), Inches(9.3), Inches(1.8))
    box.fill.solid(); box.fill.fore_color.rgb = C_PRI; box.line.color.rgb = C_ACCENT
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"  评级：增持（Overweight）"; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = C_WHITE
    p2 = tf.add_paragraph()
    if quote:
        p2.text = f"  当前价：¥{quote['price']}  |  PE(TTM): {quote.get('pe','N/A')}x  |  PB: {quote.get('pb','N/A')}x"
    p2.font.size = Pt(18); p2.font.color.rgb = C_ACCENT; p2.space_before = Pt(10)
    p3 = tf.add_paragraph(); p3.text = "  目标价与上行空间：详见估值分析页（基于行业可比PE）"; p3.font.size = Pt(14); p3.font.color.rgb = C_WHITE; p3.space_before = Pt(8)
    _add_bullets(s, [
        f"数据支撑：{len(fin)}年财务数据 + {len(hist)}条K线 + {len(peers)}家可比公司",
        "估值方法：PE比较法 / 足球场估值区间",
        "操作建议：建议结合最新市场动态和公司公告做出投资决策",
        "风险提示：投资有风险，入市需谨慎，过往业绩不代表未来表现",
    ], Inches(1), Inches(3.3), Inches(11), Inches(3), size=14)

    # ── S12: Disclaimer ──
    s = blank(); _add_header(s, "免责声明  |  Disclaimer")
    disc = f"""
本报告由研究团队编制，仅供内部参考，不构成任何投资建议。

数据来源：
  • 实时行情：腾讯财经 (qt.gtimg.cn) — 实时报价数据
  • 历史K线：腾讯财经 (web.ifzq.gtimg.cn) — 前复权日线数据 ({len(hist)}条)
  • 财务摘要：同花顺THS via AkShare — {len(fin)}年年度财务数据
  • 可比公司：腾讯财经批量接口 — {len(peers)}家实时估值

报告中的信息来源于公开渠道，不保证其完整性和准确性。
投资有风险，入市需谨慎。投资者应独立做出投资决策，并承担相应风险。
过往业绩不代表未来表现。

本报告版权归编制机构所有，未经许可不得对外披露或传播。
报告日期：{today}  |  数据存在延迟，以交易所公告为准

生成工具：china-pptx-author Skill (Generic A-share Pitch Deck Generator)
"""
    _add_text(s, disc.strip(), Inches(1), Inches(1.3), Inches(11), Inches(5), size=11, color=C_GRAY)

    return prs


# ═══════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════

def main():
    global CHART_DIR
    parser = argparse.ArgumentParser(description="Generic A-share Pitch Deck Generator")
    parser.add_argument("--company", required=True, help="公司名称")
    parser.add_argument("--ticker", required=True, help="A股代码")
    parser.add_argument("--industry", default="", help="行业名称")
    parser.add_argument("--output", default="", help="输出路径")
    parser.add_argument("--type", default="pitch", help="deck类型")
    parser.add_argument("--title", default="", help="报告标题")
    parser.add_argument("--start", default="20240101", help="K线起始日期")
    parser.add_argument("--end", default="20250608", help="K线结束日期")
    args = parser.parse_args()

    if not args.output:
        safe = "".join(c if c.isalnum() else "_" for c in args.company)
        args.output = f"{safe}_{args.ticker}_路演PPT.pptx"

    CHART_DIR = f"/tmp/{args.ticker}_charts"
    os.makedirs(CHART_DIR, exist_ok=True)

    global quote
    quote = None

    print("=" * 60)
    print(f"  {args.company} ({args.ticker}) — 路演PPT生成")
    print("=" * 60)

    # 1. Quote
    print(f"\n[1/5] 实时行情: {args.ticker}...")
    quote = fetch_quote(args.ticker)
    if quote:
        print(f"  ✓ {quote['name']}  ¥{quote['price']}  PE:{quote.get('pe')}  PB:{quote.get('pb')}  市值:{quote.get('mcap_yi',0):.0f}亿")
    else:
        print(f"  ✗ 行情获取失败")

    # 2. Financials (同花顺)
    print(f"\n[2/5] 财务数据: {args.ticker}...")
    fin = fetch_financials(args.ticker)
    if not fin.empty:
        print(f"  ✓ {len(fin)} 年数据")
        show = [c for c in ['报告期','营业总收入','净利润','销售毛利率','销售净利率','净资产收益率'] if c in fin.columns]
        if show:
            print(fin[show].tail(5).to_string(index=False))
    else:
        print("  ⚠ 财务数据获取失败")

    # 3. K-line (Tencent)
    print(f"\n[3/5] 历史K线: {args.start} ~ {args.end}...")
    hist = fetch_kline(args.ticker, args.start, args.end)
    if not hist.empty:
        print(f"  ✓ {len(hist)} 条")
        print(f"    最高 ¥{hist['最高'].max():.2f}  最低 ¥{hist['最低'].min():.2f}  最新 ¥{hist.iloc[-1]['收盘']:.2f}")
    else:
        print("  ⚠ K线获取失败")

    # 4. Peers (Tencent batch)
    print(f"\n[4/5] 行业可比公司...")
    peers = []
    peer_tickers_map = {
        '电池': ['002594','600438','002460','002129','688981','300014','600731','300450'],
        '白酒': ['000858','000568','002304','000596','603369','603197','000860'],
        '新能源': ['002594','600438','300750','002460','688981','300014','601633','300450'],
        '半导体': ['002371','603501','688981','600745','002049','603893','688036'],
        '金融': ['601318','600036','601398','600000','601166','000001','600030'],
        '医药': ['300015','000538','600276','000963','300003','002007','600196'],
    }
    industry = args.industry or '电池'
    peer_codes = peer_tickers_map.get(industry, ['002594','600438','002129','002460'])
    # Add the target company to the list for highlight
    all_codes = [args.ticker] + peer_codes
    peers = fetch_peers(all_codes)
    # Remove target from peers list for clean comparison
    peers = [p for p in peers if p['code'] != args.ticker]
    if peers:
        print(f"  ✓ {len(peers)} 家可比公司:")
        for p in peers[:8]:
            if p['price']:
                print(f"    {p['name']} ({p['code']}): ¥{p['price']:.2f}  PE:{p['pe']}  PB:{p['pb']}")

    # 5. Charts
    print(f"\n[5/5] 生成图表...")
    charts = make_charts(fin, hist, peers, args.company, args.ticker, CHART_DIR)
    print(f"  ✓ {len(charts)} 张: {CHART_DIR}/")

    # 6. PPT
    print(f"\n{'='*60}")
    print("  生成 PowerPoint...")
    prs = build_deck(args, quote, fin, hist, peers, charts)
    prs.save(args.output)
    print(f"\n  ✅ 完成: {args.output}")
    print(f"  共 {len(prs.slides)} 页  |  图表: {CHART_DIR}/")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
